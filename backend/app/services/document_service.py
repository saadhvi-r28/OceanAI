from docx import Document
from docx.shared import Inches, Pt
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from typing import List
import io

from app.models.project import Project, Section


# Color theme definitions (background, title, text, accent)
COLOR_THEMES = {
    "blue_purple": {
        "bg1": RGBColor(25, 35, 70),      # Deep navy blue
        "bg2": RGBColor(60, 40, 100),     # Rich purple
        "title": RGBColor(255, 255, 255), # White
        "text": RGBColor(240, 242, 255),  # Light blue-white
        "accent": RGBColor(100, 150, 255) # Bright blue
    },
    "green_teal": {
        "bg1": RGBColor(10, 60, 80),      # Deep ocean teal
        "bg2": RGBColor(15, 95, 75),      # Emerald green
        "title": RGBColor(255, 255, 255), # White
        "text": RGBColor(230, 255, 240),  # Light mint
        "accent": RGBColor(100, 255, 200) # Bright teal
    },
    "orange_red": {
        "bg1": RGBColor(80, 20, 20),      # Deep crimson
        "bg2": RGBColor(120, 50, 20),     # Burnt orange
        "title": RGBColor(255, 255, 255), # White
        "text": RGBColor(255, 240, 230),  # Light peach
        "accent": RGBColor(255, 140, 60)  # Bright orange
    },
    "navy_gold": {
        "bg1": RGBColor(15, 30, 60),      # Midnight navy
        "bg2": RGBColor(40, 50, 70),      # Slate blue
        "title": RGBColor(255, 215, 100), # Gold
        "text": RGBColor(255, 255, 255),  # White
        "accent": RGBColor(255, 200, 50)  # Bright gold
    },
    "pink_purple": {
        "bg1": RGBColor(80, 20, 70),      # Deep magenta
        "bg2": RGBColor(60, 30, 90),      # Royal purple
        "title": RGBColor(255, 255, 255), # White
        "text": RGBColor(255, 240, 250),  # Light pink
        "accent": RGBColor(255, 100, 200) # Bright pink
    },
    "forest_sage": {
        "bg1": RGBColor(30, 50, 40),      # Deep forest green
        "bg2": RGBColor(50, 70, 60),      # Sage
        "title": RGBColor(255, 255, 255), # White
        "text": RGBColor(240, 255, 245),  # Light sage
        "accent": RGBColor(150, 255, 180) # Bright lime
    },
    "sunset": {
        "bg1": RGBColor(60, 30, 70),      # Deep violet
        "bg2": RGBColor(100, 40, 50),     # Wine red
        "title": RGBColor(255, 220, 150), # Warm gold
        "text": RGBColor(255, 245, 235),  # Light cream
        "accent": RGBColor(255, 150, 100) # Coral
    },
    "ocean": {
        "bg1": RGBColor(10, 40, 80),      # Deep ocean blue
        "bg2": RGBColor(20, 70, 100),     # Ocean blue
        "title": RGBColor(255, 255, 255), # White
        "text": RGBColor(230, 245, 255),  # Ice blue
        "accent": RGBColor(100, 200, 255) # Sky blue
    },
    "slate": {
        "bg1": RGBColor(40, 45, 50),      # Charcoal
        "bg2": RGBColor(60, 65, 75),      # Slate gray
        "title": RGBColor(255, 255, 255), # White
        "text": RGBColor(240, 242, 245),  # Light gray
        "accent": RGBColor(150, 200, 255) # Light blue
    },
}


class DocumentService:
    @staticmethod
    def create_docx(project: Project, sections: List[Section]) -> io.BytesIO:
        """Create a Word document from project data."""
        doc = Document()
        
        # Add title
        title = doc.add_heading(project.title, 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # Add description if exists
        if project.description:
            desc_para = doc.add_paragraph(project.description)
            desc_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            doc.add_paragraph()  # Empty line
        
        # Add topic information
        topic_para = doc.add_paragraph()
        topic_para.add_run("Topic: ").bold = True
        topic_para.add_run(project.topic)
        topic_para.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
        doc.add_paragraph()  # Empty line
        
        # Sort sections by order
        sorted_sections = sorted(sections, key=lambda x: x.order)
        
        # Add each section
        for section in sorted_sections:
            # Add section heading
            doc.add_heading(section.title, level=1)
            
            # Add section content
            if section.content:
                content_para = doc.add_paragraph(section.content)
                content_para.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
            else:
                placeholder_para = doc.add_paragraph("[Content not yet generated]")
                placeholder_para.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
            
            doc.add_paragraph()  # Empty line between sections
        
        # Save to BytesIO
        file_stream = io.BytesIO()
        doc.save(file_stream)
        file_stream.seek(0)
        return file_stream
    
    @staticmethod
    def create_pptx(project: Project, sections: List[Section]) -> io.BytesIO:
        """Create a PowerPoint presentation from project data."""
        prs = Presentation()
        prs.slide_width = PptxInches(12)
        prs.slide_height = PptxInches(8.5)
        
        # Get color theme
        theme = COLOR_THEMES.get(project.color_theme, COLOR_THEMES["blue_purple"])
        
        # Title slide with custom design
        blank_layout = prs.slide_layouts[6]  # Blank layout
        title_slide = prs.slides.add_slide(blank_layout)
        
        # Add gradient background to title slide
        background = title_slide.background
        fill = background.fill
        fill.gradient()
        fill.gradient_angle = 45
        fill.gradient_stops[0].color.rgb = theme["bg1"]
        fill.gradient_stops[1].color.rgb = theme["bg2"]
        
        # Add title text box - centered and full width
        title_box = title_slide.shapes.add_textbox(
            PptxInches(0.5), 
            PptxInches(3), 
            PptxInches(11), 
            PptxInches(1.5)
        )
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_frame.vertical_anchor = 1  # Middle vertical alignment
        title_para = title_frame.paragraphs[0]
        title_para.text = project.title
        title_para.alignment = PP_ALIGN.CENTER
        title_para.font.size = PptxPt(48)
        title_para.font.bold = True
        title_para.font.color.rgb = theme["title"]
        
        # Add subtitle/description text box - full width
        subtitle_box = title_slide.shapes.add_textbox(
            PptxInches(1), 
            PptxInches(5), 
            PptxInches(10), 
            PptxInches(2)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.word_wrap = True
        subtitle_frame.vertical_anchor = 1  # Middle vertical alignment
        subtitle_para = subtitle_frame.paragraphs[0]
        if project.description:
            subtitle_para.text = project.description
        else:
            subtitle_para.text = project.topic
        subtitle_para.alignment = PP_ALIGN.CENTER
        subtitle_para.font.size = PptxPt(24)
        subtitle_para.font.color.rgb = theme["text"]
        
        # Sort sections by order
        sorted_sections = sorted(sections, key=lambda x: x.order)
        
        # Add content slides
        for section in sorted_sections:
            # Use blank layout for custom positioning
            blank_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(blank_layout)
            
            # Add gradient background with theme colors
            background = slide.background
            fill = background.fill
            fill.gradient()
            fill.gradient_angle = 45
            fill.gradient_stops[0].color.rgb = theme["bg1"]
            fill.gradient_stops[1].color.rgb = theme["bg2"]
            
            # Add title at the top - centered with styling
            title_box = slide.shapes.add_textbox(
                PptxInches(0.5), 
                PptxInches(0.5), 
                PptxInches(11), 
                PptxInches(1.2)
            )
            title_frame = title_box.text_frame
            title_frame.word_wrap = True
            title_para = title_frame.paragraphs[0]
            title_para.text = section.title
            title_para.alignment = PP_ALIGN.CENTER
            title_para.font.size = PptxPt(32)
            title_para.font.bold = True
            title_para.font.color.rgb = theme["title"]
            
            # Add content box - centered on the page both horizontally and vertically
            # Slide is 10" wide x 8.5" tall
            # Box is 7" wide x 4.5" tall, so:
            # Left = (10 - 7) / 2 = 1.5"
            # Top = (8.5 - 4.5) / 2 = 2"
            content_box = slide.shapes.add_textbox(
                PptxInches(1.75), 
                PptxInches(2), 
                PptxInches(7), 
                PptxInches(4.5)
            )
            text_frame = content_box.text_frame
            text_frame.vertical_anchor = 1  # MSO_ANCHOR.MIDDLE for vertical centering
            text_frame.word_wrap = True
            
            if section.content:
                # Remove redundant bullet markers (• or - at start) and split content
                content_lines = section.content.strip().split('\n')
                for i, line in enumerate(content_lines):
                    line = line.strip()
                    # Remove leading bullet characters but keep track if it was bulleted
                    is_bullet = False
                    if line.startswith('• '):
                        line = line[2:]
                        is_bullet = True
                    elif line.startswith('- '):
                        line = line[2:]
                        is_bullet = True
                    elif line.startswith('* '):
                        line = line[2:]
                        is_bullet = True
                    
                    if line:
                        if i == 0:
                            p = text_frame.paragraphs[0]
                        else:
                            p = text_frame.add_paragraph()
                        
                        # Add bullet symbol back if it was originally bulleted
                        if is_bullet:
                            p.text = '• ' + line
                        else:
                            p.text = line
                        
                        p.level = 0
                        p.alignment = PP_ALIGN.JUSTIFY  # Justify text
                        p.font.size = PptxPt(20)
                        p.font.color.rgb = theme["text"]
                        p.space_before = PptxPt(8)
                        p.space_after = PptxPt(8)
            else:
                p = text_frame.paragraphs[0]
                p.text = "[Content not yet generated]"
                p.alignment = PP_ALIGN.CENTER
                p.font.size = PptxPt(18)
                p.font.color.rgb = theme["text"]
        
        # Save to BytesIO
        file_stream = io.BytesIO()
        prs.save(file_stream)
        file_stream.seek(0)
        return file_stream


document_service = DocumentService()
